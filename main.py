import os
import io
import datetime
from telegram import Update
from telegram.ext import Application, CommandHandler, MessageHandler, ContextTypes, filters
from pptx import Presentation

# --- Konfiguratsiya va Global Sozlamalar ---
# **MUHIM: ALMASHTIRING** Oʻz Telegram bot tokeningizni joylang
BOT_TOKEN = "8579631704:AAHxcJpfN0sFC4C8N8GJHPWpLXsMe3dQ0qQ"
# **MUHIM: ALMASHTIRING** Sizning Telegram ID raqamingiz (son)
ADMIN_USER_ID = 8005357331  

# --- Global Holatni Boshqarish ---
user_states = {} # Foydalanuvchi holatini saqlash: awaiting_topic, awaiting_pptx, awaiting_content
user_data = {}  # Foydalanuvchining mavzusi va fayl ma'lumotlarini saqlash

# --- PPTX Matnini Almashtirish Funksiyasi (Xato Tuzatildi) ---

def replace_text_in_slides(prs, new_texts_list):
    """
    Taqdimotdagi matn qutilaridagi matnni yangi matnlar bilan almashtiradi.
    Matnni almashtirishning eng xavfsiz usuli qo'llanildi.
    """
    text_index = 0
    
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                
                # Slayd ichidagi barcha paragraflarni aylanib chiqamiz
                for paragraph in text_frame.paragraphs:
                    if text_index < len(new_texts_list):
                        
                        # Faqat birinchi Run obyektining matnini almashtiramiz
                        if paragraph.runs:
                            # Mavjud barcha matnlarni tozalab, birinchi run'ga yangi matnni yozamiz
                            paragraph.runs[0].text = new_texts_list[text_index].strip()
                            # Qo'shimcha run larni o'chiramiz (agar mavjud bo'lsa)
                            for i in range(len(paragraph.runs) - 1, 0, -1):
                                paragraph._p.remove(paragraph.runs[i]._r)
                        else:
                            # Agar paragraf bo'sh bo'lsa, yangi run qo'shish
                            new_run = paragraph.add_run()
                            new_run.text = new_texts_list[text_index].strip()
                        
                        text_index += 1
                    
                    if text_index >= len(new_texts_list):
                        return 

# --- Telegram Bot Handlerlari ---

async def start_command(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    """/start buyrug'ini qabul qiladi."""
    user_id = update.effective_user.id
    today = datetime.date.today()
    target_date = datetime.date(2025, 1, 1)

    # 1-yanvargacha boshqa foydalanuvchilar uchun cheklov
    if user_id != ADMIN_USER_ID:
        if today < target_date:
            await update.message.reply_text(
                f"Assalomu alaykum! 👋\n\n"
                f"Men WPS Office shablonlarini avtomatik to'ldiruvchi botman. **2025 yil 1 yanvardan** to'liq ishga tushaman! Kutib qoling! 😉"
            )
            return

    # Admin va 1-yanvardan keyingi foydalanuvchilar uchun asosiy jarayon
    await update.message.reply_text(
        f"Salom! Ishni boshlash uchun:\n\n"
        f"1. Prezentatsiya **mavzusini** yozing."
    )
    user_states[user_id] = 'awaiting_topic'
    # Oldingi ma'lumotlarni tozalash
    user_data[user_id] = {} 

async def handle_message(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    """Matn yoki fayllarni qabul qiladi va jarayonni boshqaradi."""
    user_id = update.effective_user.id
    state = user_states.get(user_id)
    
    # Reklama davrida admin bo'lmaganlarga xabar bermaslik
    today = datetime.date.today()
    target_date = datetime.date(2025, 1, 1)
    if user_id != ADMIN_USER_ID and today < target_date:
        return

    # 1. Mavzuni qabul qilish
    if state == 'awaiting_topic' and update.message.text:
        user_data[user_id]['topic'] = update.message.text
        user_states[user_id] = 'awaiting_pptx'
        await update.message.reply_text(
            f"A'lo! Mavzu: **{update.message.text}**.\n\n"
            f"2. Endi iltimos, WPS/PowerPoint'dan tanlagan **PPTX shabloningizni** fayl sifatida yuboring."
        )

    # 2. PPTX faylini qabul qilish
    elif state == 'awaiting_pptx' and update.message.document and update.message.document.file_name.endswith('.pptx'):
        await update.message.reply_text("Fayl qabul qilindi. Kontent uchun prompt yaratilmoqda...")
        
        # Fayl ma'lumotlarini vaqtinchalik xotiraga saqlash
        pptx_file = await update.message.document.get_file()
        file_id = update.message.document.file_id
        
        user_data[user_id]['file_id'] = file_id
        user_data[user_id]['file_name'] = update.message.document.file_name

        # Faylni xotiraga yuklab, slaydlar sonini aniqlash
        file_data = io.BytesIO()
        await pptx_file.download_to_memory(file_data)
        file_data.seek(0)
        
        try:
            prs = Presentation(file_data)
            num_slides = len(prs.slides)
        except Exception:
            await update.message.reply_text("PPTX faylini o'qishda xato. Format to'g'riligini tekshiring.")
            user_states[user_id] = None
            return

        # Prompt yaratish (Foydalanuvchiga beriladigan ko'rsatma)
        # Slaydlar soni * 2 (Sarlavha va asosiy matn uchun)
        prompt_texts_count = num_slides * 2 
        topic = user_data[user_id]['topic']
        
        prompt = (
            f"Quyidagi mavzu bo'yicha **{prompt_texts_count} qismdan** iborat prezentatsiya matnini tayyorlang:\n"
            f"**Mavzu:** {topic}\n\n"
            f"**Talablar:**\n"
            f"1. Har bir qism alohida qatorda (yoki raqamlanib) bo'lsin.\n"
            f"2. Har bir qism 1-4 jumla atrofida bo'lsin (Qisqa va aniq).\n"
            f"3. Ohang: Ilmiy-rasmiy (o'zbek tilida, professional).\n\n"
            f"**Natijani to'g'ridan-to'g'ri faqat matn qismlari bilan yuboring, boshqa kirish matnlarini qo'shmang!**"
        )
        
        user_states[user_id] = 'awaiting_content'
        await update.message.reply_text(
            f"3. **Mavzuga mos kontentni olish uchun**:\n\n"
            f"Sizga kerak bo'ladigan **Prompt** (ko'rsatma):\n\n"
            f"```\n{prompt}\n```\n\n"
            f"Ushbu promptni istalgan AI ga (Gemini, ChatGPT) bering. Olingan toza matnni menga **javob sifatida yuboring** (Reply)."
        )

    # 3. Kontentni qabul qilish va faylni to'ldirish
    elif state == 'awaiting_content' and update.message.text:
        await update.message.reply_text("Kontent qabul qilindi. Matnni PPTXga joylamoqdaman...")
        
        # AI dan kelgan matnni tozalash
        raw_content = update.message.text
        # Har bir qatorni alohida matn qismi deb hisoblaymiz
        new_texts = [text.strip() for text in raw_content.split('\n') if text.strip()]

        file_id = user_data[user_id]['file_id']
        file_name = user_data[user_id]['file_name']
        topic = user_data[user_id]['topic']

        # Faylni qayta yuklash
        pptx_file = await context.bot.get_file(file_id)
        file_data = io.BytesIO()
        await pptx_file.download_to_memory(file_data)
        file_data.seek(0)
        
        try:
            prs = Presentation(file_data)
            # 4. Matnni joylashtirish
            replace_text_in_slides(prs, new_texts)

            # 5. Yangi faylni saqlash va yuborish
            output_buffer = io.BytesIO()
            prs.save(output_buffer)
            output_buffer.seek(0)

            await update.message.reply_document(
                document=output_buffer,
                filename=f"To'ldirilgan_{file_name}",
                caption=f"Tayyor prezentatsiya:\n\n**Mavzu:** {topic}"
            )
        except Exception as e:
            # Xatolikni admin uchun ko'rsatish
            await update.message.reply_text(f"Faylga matn joylashda kutilmagan xato yuz berdi: {e}")

        # Jarayonni tugatish
        user_states[user_id] = None
        if user_id in user_data:
            del user_data[user_id]

    else:
        # Noto'g'ri turdagi xabar
        await update.message.reply_text("Noto'g'ri qadam. Iltimos, /start buyrug'ini qaytadan bosing.")


def main() -> None:
    """Botni ishga tushiradi."""
    application = Application.builder().token(BOT_TOKEN).build()

    application.add_handler(CommandHandler("start", start_command))
    application.add_handler(MessageHandler(filters.TEXT | filters.Document.ALL, handle_message))

    application.run_polling(allowed_updates=Update.ALL_TYPES)

if __name__ == '__main__':
    main()
    
